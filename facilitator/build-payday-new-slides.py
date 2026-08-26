#!/usr/bin/env python3
"""
Build the five new Payday slides, move 132 after 80, and keep the six hidden slides hidden.

Reads   : AI in Testing Workshop - Payday 26Aug2026.pptx   (never modified)
Writes  : AI in Testing Workshop - Payday 26Aug2026 - FINAL.pptx
          handout-slide-101.pdf                             (via LibreOffice, separately)

Run it from the folder that holds the deck:

    python3 build-payday-new-slides.py

Everything is idempotent in the sense that it always builds from the untouched input, so
re-running gives the same output. It is safe to run as many times as you like.

------------------------------------------------------------------------------------------
Design notes, so the slides do not look bolted on
------------------------------------------------------------------------------------------
The deck is Google-Slides-authored on the `simple-light-2` master. Measured from the
existing slides rather than guessed:

    title      placeholder at (0.59", 0.59"), 8.83" wide, colour #374151
    body       Open Sans / Open Sans Light, 11-12pt
    palette    accent1 #4285F4, dk2 #595959, lt2 #EEEEEE   (theme colour scheme)
    canvas     10" x 5.625"  (16:9)

So the new slides use the master's own TITLE_ONLY / BLANK layouts, the same title
geometry, Open Sans throughout, and the theme's own greys. Tables get their PowerPoint
banding stripped and explicit hairline borders, because the default python-pptx table
style is bright blue and would stand out badly.

------------------------------------------------------------------------------------------
Why NEW-4 is two slides
------------------------------------------------------------------------------------------
The spec asks for a click-reveal. python-pptx cannot author animations at all, and
hand-writing the animation XML is fragile across PowerPoint, Keynote and Google Slides.
Two consecutive slides give the identical effect - first row, click, second row - and work
in every viewer. They are NEW-4a and NEW-4b.
"""

from __future__ import annotations

import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC = Path("AI in Testing Workshop - Payday 26Aug2026.pptx")
OUT = Path("AI in Testing Workshop - Payday 26Aug2026 - FINAL.pptx")

# Slides that must remain hidden. Tracked by slide_id, never by position, because every
# insert and move shifts positions underneath us.
MUST_STAY_HIDDEN = [151, 152, 160, 167, 269, 270]

FONT = "Open Sans"
FONT_LIGHT = "Open Sans Light"

INK = RGBColor(0x37, 0x41, 0x51)      # title / body, matches slide 3
MUTED = RGBColor(0x59, 0x59, 0x59)    # theme dk2
FAINT = RGBColor(0x80, 0x86, 0x8B)
ACCENT = RGBColor(0x42, 0x85, 0xF4)   # theme accent1
RULE = RGBColor(0xDA, 0xDC, 0xE0)     # hairline borders
HEAD_BG = RGBColor(0xF1, 0xF3, 0xF4)  # header row fill, a shade off theme lt2
BAD = RGBColor(0xB3, 0x14, 0x12)      # the one number that costs money

TITLE_POS = (Inches(0.59), Inches(0.59), Inches(8.83), Inches(0.63))


# ---------------------------------------------------------------------------
# primitives
# ---------------------------------------------------------------------------

def layout(prs: Presentation, name: str):
    for lay in prs.slide_masters[0].slide_layouts:
        if lay.name == name:
            return lay
    sys.exit(f"layout {name!r} not found on master 0")


def add_slide(prs: Presentation, layout_name: str = "BLANK"):
    """Append a slide and strip the inherited slide-number placeholder."""
    slide = prs.slides.add_slide(layout(prs, layout_name))
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx == 12:
            shape._element.getparent().remove(shape._element)
    return slide


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def write(frame, blocks):
    """Fill a text frame. Each block is (runs, opts).

    runs  - a string, or a list of (text, {run overrides}) tuples
    opts  - size, bold, colour, font, space_before, space_after, line, align
    """
    frame.clear()
    for i, (runs, opts) in enumerate(blocks):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        if opts.get("space_before"):
            para.space_before = Pt(opts["space_before"])
        if opts.get("space_after"):
            para.space_after = Pt(opts["space_after"])
        if opts.get("line"):
            para.line_spacing = opts["line"]
        if opts.get("align"):
            para.alignment = opts["align"]
        if isinstance(runs, str):
            runs = [(runs, {})]
        for text, over in runs:
            run = para.add_run()
            run.text = text
            font = run.font
            font.name = over.get("font", opts.get("font", FONT))
            font.size = Pt(over.get("size", opts.get("size", 12)))
            font.bold = over.get("bold", opts.get("bold", False))
            font.italic = over.get("italic", opts.get("italic", False))
            font.color.rgb = over.get("colour", opts.get("colour", INK))
    return frame


def title(slide, text, size=20):
    frame = textbox(slide, *TITLE_POS)
    frame.word_wrap = True
    write(frame, [(text, {"size": size, "bold": True, "colour": INK})])
    return frame


def cell_border(cell, colour=RULE, width=Pt(0.75)):
    """python-pptx has no border API; write the four lines into tcPr directly."""
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for existing in tc_pr.findall(qn(tag)):
            tc_pr.remove(existing)
    # order matters in the schema: lnL, lnR, lnT, lnB come before the fill
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):
        ln = tc_pr.makeelement(qn(tag), {"w": str(int(width)), "cap": "flat"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": f"{colour}"})
        fill.append(clr)
        ln.append(fill)
        tc_pr.insert(0, ln)


def make_table(slide, rows, left, top, width, col_widths, *, size=9,
               header=True, row_height=Pt(20)):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                   Inches(0.3 * n_rows))
    table = shape.table

    # Strip the default blue table style and its banding.
    tbl_pr = table._tbl.find(qn("a:tblPr"))
    if tbl_pr is not None:
        for attr in ("firstRow", "firstCol", "bandRow", "bandCol", "lastRow", "lastCol"):
            if attr in tbl_pr.attrib:
                del tbl_pr.attrib[attr]
        for style_id in tbl_pr.findall(qn("a:tableStyleId")):
            tbl_pr.remove(style_id)

    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for row in table.rows:
        row.height = row_height

    for r, row in enumerate(rows):
        for c, spec in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell_border(cell)
            is_head = header and r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEAD_BG if is_head else RGBColor(0xFF, 0xFF, 0xFF)
            opts = {
                "size": size,
                "bold": is_head,
                "colour": MUTED if is_head else INK,
                "font": FONT,
            }
            if isinstance(spec, tuple):
                text, over = spec
                opts.update(over)
            else:
                text = spec
            write(cell.text_frame, [(text, opts)])
    return table


def footnote(slide, text, top=Inches(4.95), colour=FAINT, size=9, italic=True):
    frame = textbox(slide, Inches(0.59), top, Inches(8.83), Inches(0.4))
    write(frame, [(text, {"size": size, "colour": colour, "italic": italic,
                          "font": FONT_LIGHT})])
    return frame


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# reordering, by slide_id so positions cannot bite us
# ---------------------------------------------------------------------------

def index_of(prs: Presentation, slide_id: int) -> int:
    for i, slide in enumerate(prs.slides):
        if slide.slide_id == slide_id:
            return i
    sys.exit(f"slide_id {slide_id} vanished")


def move_after(prs: Presentation, slide_id: int, anchor_id: int) -> None:
    """Move the slide with slide_id so it sits directly after anchor_id."""
    lst = prs.slides._sldIdLst
    entries = list(lst)
    moving = entries[index_of(prs, slide_id)]
    lst.remove(moving)
    entries = list(lst)                       # re-read: indices shifted
    anchor_pos = next(i for i, e in enumerate(entries)
                      if int(e.get("id")) == anchor_id)
    lst.insert(anchor_pos + 1, moving)


def set_hidden(slide, hidden: bool) -> None:
    element = slide._element
    if hidden:
        element.set("show", "0")
    elif "show" in element.attrib:
        del element.attrib["show"]


# ---------------------------------------------------------------------------
# the five slides
# ---------------------------------------------------------------------------

AGENDA = [
    ("Reykjavik", "Bucharest", "Session"),
    ("08:30", "11:30", "Welcome, and why we test"),
    ("09:15", "12:15", "Context and prompting essentials"),
    ("09:40", "12:40", "Break"),
    ("09:50", "12:50", "SDLC I  -  Discovery and requirement analysis"),
    ("10:20", "13:20", "SDLC II  -  Test planning and design"),
    ("11:00", "14:00", "Lunch"),
    ("12:00", "15:00", "SDLC III  -  From zero to a green e2e suite"),
    ("12:50", "15:50", "Break"),
    ("13:00", "16:00", "The customization toolbox: rules, prompts, skills, agents, hooks"),
    ("13:50", "16:50", "Playwright agents, the pipeline, and one very expensive krona"),
    ("14:30", "17:30", "Release, defects, and running it without you"),
    ("14:45", "17:45", "What we saw, what stays yours, what you do next week"),
]

TOOLING = [
    ("", "Read by both", "Cursor", "GitHub Copilot"),
    ("Always-on instructions", "AGENTS.md",
     ".cursor/rules/*.mdc  (alwaysApply)", ".github/copilot-instructions.md"),
    ("Path-scoped instructions", "-",
     ".mdc with  globs:", "*.instructions.md with  applyTo:"),
    ("Reusable prompts", ".agents/skills/*/SKILL.md",
     "skills, slash-invocable   (.cursor/commands = legacy)",
     ".github/prompts/*.prompt.md"),
    ("Subagents", "-",
     ".cursor/agents/*.md  (readonly:)   run with  /name",
     ".github/agents/*.agent.md  (tools:)   ask in prose"),
    ("Agent handoffs", "-",
     "no field  -  orchestrator pattern",
     "handoffs:  renders a button   (not in cloud agents)"),
    ("Hooks (can block a write)", "-",
     ".cursor/hooks.json", ".github/hooks/*.json"),
    ("MCP servers", "-",
     ".cursor/mcp.json  ->  mcpServers", ".vscode/mcp.json  ->  servers"),
]


def build_new_1(prs):
    slide = add_slide(prs, "BLANK")
    title(slide, "Agenda")
    make_table(
        slide, AGENDA,
        left=Inches(0.59), top=Inches(1.32), width=Inches(8.83),
        col_widths=[Inches(1.0), Inches(1.05), Inches(6.78)],
        size=9.5, row_height=Pt(18),
    )
    footnote(slide, "One hour for lunch, a break each side of it. We finish on time.",
             top=Inches(5.06))
    notes(slide, "Replaces the untimed agenda on the previous slide.\n\n"
                 "Two clocks on purpose - theirs first, because the room is in Reykjavik "
                 "and you are three hours ahead in Bucharest.\n\n"
                 "Leave 'one very expensive krona' unexplained. It seeds the 16:50 reveal "
                 "all day and somebody will ask about it at lunch. If they ask, say: "
                 "'you'll see - and you'll be able to tell me exactly how expensive.'")
    return slide


def build_new_2(prs):
    slide = add_slide(prs, "BLANK")
    title(slide, "Payday salary run  -  the demo app")
    frame = textbox(slide, Inches(0.59), Inches(1.42), Inches(5.5), Inches(2.6))
    write(frame, [
        ([("Add employees", {"bold": True}),
          ("   -   name, kennitala, monthly gross in ISK", {})],
         {"size": 13, "space_after": 9}),
        ([("Runs the numbers", {"bold": True}),
          ("   -   pension, tax, net, paid state, totals", {})],
         {"size": 13, "space_after": 9}),
        ([("Three files", {"bold": True}),
          ("   -   ", {}),
          ("index.html", {"colour": ACCENT}), ("  ", {}),
          ("main.js", {"colour": ACCENT}), ("  ", {}),
          ("style.css", {"colour": ACCENT})],
         {"size": 13, "space_after": 9}),
        ([("No framework, no backend, state in the browser.", {})],
         {"size": 12, "colour": MUTED, "font": FONT_LIGHT}),
    ])

    frame = textbox(slide, Inches(0.59), Inches(3.25), Inches(8.83), Inches(1.2))
    write(frame, [
        ([("It has bugs in it on purpose.", {})],
         {"size": 22, "bold": True, "colour": INK, "space_after": 6}),
        ([("Some of you will find them before I show you. When you do, write it down "
           "and don't fix it.", {})],
         {"size": 13, "colour": MUTED, "font": FONT_LIGHT}),
    ])

    frame = textbox(slide, Inches(6.25), Inches(1.42), Inches(3.16), Inches(2.0))
    write(frame, [
        ([("Deliberately trivial", {"bold": True})],
         {"size": 11, "colour": MUTED, "space_after": 5}),
        ([("so every demo takes a minute instead of twenty. Every technique today "
           "scales to your real product - the app is just a fast target.", {})],
         {"size": 11, "colour": MUTED, "font": FONT_LIGHT}),
    ])
    notes(slide, "DO NOT SAY HOW MANY BUGS. Ten are planted; nine are verified "
                 "reproducible; the tenth is a crash.\n\n"
                 "The 'write it down and don't fix it' instruction matters - it keeps "
                 "findings available for the reveals later, and it models the gate.")
    return slide


def build_new_3(prs):
    slide = add_slide(prs, "BLANK")
    title(slide, "One repo, two editors")
    make_table(
        slide, TOOLING,
        left=Inches(0.59), top=Inches(1.30), width=Inches(8.83),
        col_widths=[Inches(1.60), Inches(1.72), Inches(2.70), Inches(2.81)],
        size=8, row_height=Pt(17),
    )
    frame = textbox(slide, Inches(0.59), Inches(3.52), Inches(8.83), Inches(0.85))
    write(frame, [
        ([("AGENTS.md", {"colour": ACCENT, "bold": True}),
          (" and ", {}),
          (".agents/skills/", {"colour": ACCENT, "bold": True}),
          (" are read by both - put shared knowledge there and maintain it once.", {})],
         {"size": 10.5, "colour": INK, "space_after": 3}),
        ([("The MCP key differs: ", {}),
          ("mcpServers", {"bold": True}), (" in Cursor, ", {}),
          ("servers", {"bold": True}),
          (" in VS Code. Copying the file across without changing the key is the most "
           "common setup failure in this whole space.", {})],
         {"size": 10.5, "colour": INK, "space_after": 3}),
        ([("Both tools have hooks that can block a tool call. This is not a "
           "Claude-only trick.", {})],
         {"size": 10.5, "colour": INK}),
    ])
    notes(slide, "Three lines to say over this, in order:\n\n"
                 "1. AGENTS.md and .agents/skills/ are read by BOTH. That is the answer "
                 "to 'what happens when the next hire uses a different editor'.\n\n"
                 "2. The MCP key differs - mcpServers vs servers. Most common setup "
                 "failure in this space.\n\n"
                 "3. Both have blocking hooks.\n\n"
                 "On the handoffs row: the field is real in VS Code / Copilot and renders "
                 "as a button after a response. Cursor has no such field, and Copilot's "
                 "cloud agents on github.com ignore it on purpose. So the repo declares "
                 "the chain once and the generator writes it two ways - a handoffs: list "
                 "for Copilot, a 'Next step' section in the body for Cursor's parent "
                 "orchestrator to read. Show agents/_generate.py if they ask.\n\n"
                 "INVOCATION, because someone will ask and the wrong answer wastes "
                 "their afternoon: Cursor documents /name for a specific subagent - "
                 "/bug-hunter, /test-planner. Copilot has no slash and no @ for agents; "
                 "you ask in prose or let the model delegate. There is NO @agent syntax "
                 "in either tool. Skills are different again: they fire automatically on "
                 "a description match in both tools, and Cursor also lets you pick one "
                 "from /.")
    return slide


CLIFF_HEAD = ("Monthly gross", "Income tax", "Net pay")
CLIFF_LOW = ("468.749 kr.", "73.750 kr.", "376.249 kr.")
CLIFF_HIGH = ("468.750 kr.", "103.000 kr.", "347.000 kr.")


def cliff_slide(prs, *, revealed: bool):
    slide = add_slide(prs, "BLANK")
    title(slide, "One krona.", size=34)

    rows = [CLIFF_HEAD, CLIFF_LOW]
    if revealed:
        rows.append([(v, {"bold": True, "colour": BAD}) for v in CLIFF_HIGH])
    make_table(
        slide, rows,
        left=Inches(0.59), top=Inches(1.72), width=Inches(7.40),
        col_widths=[Inches(2.47), Inches(2.47), Inches(2.46)],
        size=15, row_height=Pt(34),
    )

    if revealed:
        frame = textbox(slide, Inches(0.59), Inches(3.45), Inches(8.83), Inches(1.3))
        write(frame, [
            ([("One krona more gross.  ", {}),
              ("29.249 kronur", {"bold": True, "colour": BAD}),
              (" less in their pocket.", {})],
             {"size": 19, "colour": INK, "space_after": 10}),
            ([("Every test in this room is green.", {})],
             {"size": 15, "colour": MUTED, "font": FONT_LIGHT, "italic": True}),
        ])
        notes(slide, "THE REVEAL. Say the number, then stop talking. Let the silence do "
                     "the work.\n\n"
                     "Mechanism: the band-2 rate is applied to the WHOLE taxable base "
                     "once the threshold is reached, instead of only the portion above "
                     "it. Taxable base is gross x 0.96, so a gross of 468.750 lands the "
                     "base exactly on the 450.000 threshold.\n\n"
                     "Then the three points, from WORKSHOP-SCRIPT-PAYDAY.md at 17:01:\n"
                     "- Invisible to any test that picks round numbers. 400.000 and "
                     "500.000 both look fine.\n"
                     "- Only found by testing both sides of a boundary, one unit apart - "
                     "which AI does well and humans skip under time pressure.\n"
                     "- Nothing in the code says a tax cliff is wrong. That is a rule "
                     "about the world, and a person has to know it.\n\n"
                     "Ask: would your current suite have caught this? Then: would you "
                     "have thought to ask for it, before an AI generated forty boundary "
                     "cases for free?")
    else:
        footnote(slide, "One employee. One krona more gross next month.",
                 top=Inches(2.95), colour=MUTED, size=15, italic=False)
        notes(slide, "Put this up and read the row out loud. Do not click yet.\n\n"
                     "Ask the room what they expect the next row to look like if the "
                     "employee earns ONE krona more. Let somebody say 'about the same'. "
                     "Then click.\n\n"
                     "The pause before the click is the slide. Give it four or five "
                     "seconds of silence - longer than is comfortable.")
    return slide


def build_new_5(prs):
    slide = add_slide(prs, "BLANK")
    title(slide, "One workflow. One gate.", size=26)
    frame = textbox(slide, Inches(0.59), Inches(1.70), Inches(8.83), Inches(2.4))
    write(frame, [
        ([("Each of you, out loud:", {})],
         {"size": 14, "colour": MUTED, "font": FONT_LIGHT, "space_after": 14}),
        ([("1.  ", {"colour": ACCENT, "bold": True}),
          ("The one workflow from today you will set up next week.", {})],
         {"size": 18, "space_after": 12}),
        ([("2.  ", {"colour": ACCENT, "bold": True}),
          ("Where the human gate sits in it.", {})],
         {"size": 18, "space_after": 12}),
    ])
    frame = textbox(slide, Inches(0.59), Inches(3.50), Inches(8.83), Inches(0.8))
    write(frame, [
        ([("Write them in the channel. That list is what this day was for.", {})],
         {"size": 14, "colour": MUTED, "font": FONT_LIGHT, "italic": True}),
    ])
    notes(slide, "Go round the room one at a time - with two to five people this works, "
                 "and it is the whole point of the day.\n\n"
                 "One workflow, not five. A person who commits to one thing does it; a "
                 "person who commits to five does none.\n\n"
                 "If somebody names a workflow with no gate in it, that is the most "
                 "useful conversation of the afternoon. Do not let it pass.\n\n"
                 "Check their answers against the whiteboard from slide 4 - what they "
                 "said they wanted this morning.")
    return slide


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------

def main() -> None:
    if not SRC.exists():
        sys.exit(f"cannot find {SRC} - run this from the folder that holds the deck")

    prs = Presentation(str(SRC))
    before = len(prs.slides)
    print(f"loaded {SRC.name}: {before} slides, "
          f"{sum(1 for s in prs.slides if s._element.get('show') != '0')} visible")

    # Capture identities BEFORE anything moves.
    ids = {n: prs.slides[n - 1].slide_id
           for n in (3, 80, 101, 126, 132, 145, *MUST_STAY_HIDDEN)}
    hidden_ids = {n: ids[n] for n in MUST_STAY_HIDDEN}
    for n, sid in hidden_ids.items():
        if prs.slides[index_of(prs, sid)]._element.get("show") != "0":
            sys.exit(f"slide {n} was expected to be hidden in the input and is not")

    # Build the new slides. They land at the end; we place them next.
    new_1 = build_new_1(prs)
    new_2 = build_new_2(prs)
    new_3 = build_new_3(prs)
    new_4a = cliff_slide(prs, revealed=False)
    new_4b = cliff_slide(prs, revealed=True)
    new_5 = build_new_5(prs)
    print(f"built 6 slides (NEW-4 is a two-slide click reveal)")

    # Place them, and move 132. Order matters: each call re-reads positions.
    move_after(prs, new_1.slide_id, ids[3])          # NEW-1 straight after the old agenda
    set_hidden(prs.slides[index_of(prs, ids[3])], True)   # ...and retire the old one

    move_after(prs, ids[132], ids[80])               # 132 follows 80, as asked
    move_after(prs, new_2.slide_id, ids[132])         # then the app
    move_after(prs, new_3.slide_id, new_2.slide_id)   # then the two dialects

    move_after(prs, new_4a.slide_id, ids[126])        # the cliff, after the lifecycle
    move_after(prs, new_4b.slide_id, new_4a.slide_id)

    move_after(prs, new_5.slide_id, ids[145])         # commitments, before the timeline

    # The whole point of tracking by id: prove the six are still hidden.
    for n, sid in hidden_ids.items():
        slide = prs.slides[index_of(prs, sid)]
        if slide._element.get("show") != "0":
            sys.exit(f"FAILED: slide originally at {n} is no longer hidden")

    prs.save(str(OUT))

    check = Presentation(str(OUT))
    after = len(check.slides)
    visible = sum(1 for s in check.slides if s._element.get("show") != "0")
    print(f"\nwrote {OUT.name}: {after} slides ({after - before} added), {visible} visible")

    print("\nnew positions:")
    for label, sid in [("NEW-1 agenda", new_1.slide_id),
                       ("132 zero-to-green", ids[132]),
                       ("NEW-2 the app", new_2.slide_id),
                       ("NEW-3 two editors", new_3.slide_id),
                       ("NEW-4a cliff, before", new_4a.slide_id),
                       ("NEW-4b cliff, revealed", new_4b.slide_id),
                       ("NEW-5 commitments", new_5.slide_id),
                       ("old agenda, now hidden", ids[3])]:
        print(f"  slide {index_of(check, sid) + 1:4}  {label}")

    print("\nhidden, as required: " + ", ".join(
        f"{n}->{index_of(check, sid) + 1}" for n, sid in hidden_ids.items()))
    print(f"handout source: slide 101 is now at {index_of(check, ids[101]) + 1}")


if __name__ == "__main__":
    main()
