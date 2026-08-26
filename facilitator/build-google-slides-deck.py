#!/usr/bin/env python3
"""
Build a Google-Slides-ready deck: the running order, and nothing else.

Google Slides discards PowerPoint custom shows on import, so the custShowLst in the FINAL
deck does nothing there. Rather than have you skip slides by hand in front of a room, this
writes a deck whose slide order IS the running order:

  * 119 slides kept, in play order
  * the 4 slides the day visits twice are physically duplicated, so it is forward-only
  * all 24 cut slides removed, along with the 133 hidden legacy slides
  * nothing to skip, nothing to hide, no numbers to hunt

Reads   : AI in Testing Workshop - Payday 26Aug2026 - FINAL.pptx
Writes  : AI in Testing Workshop - Payday 26Aug2026 - GOOGLE SLIDES.pptx

The FINAL deck is left alone - it stays the reference the playbook's slide numbers point
at, and it is what you would present from if you ever open PowerPoint proper.

Dropping the unused slide parts also takes the file from ~80MB to something Google Slides
imports without complaining.
"""

from __future__ import annotations

import copy
import importlib.util
import sys
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TARGET_MODE as RTM
from pptx.opc.package import _Relationship
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
SRC = Path("AI in Testing Workshop - Payday 26Aug2026 - FINAL.pptx")
OUT = Path("AI in Testing Workshop - Payday 26Aug2026 - GOOGLE SLIDES.pptx")


def load_path():
    for candidate in (Path("add-custom-show.py"), HERE / "add-custom-show.py"):
        if candidate.exists():
            spec = importlib.util.spec_from_file_location("acs", candidate)
            mod = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(mod)
            return mod.PATH
    sys.exit("cannot find add-custom-show.py - it holds the running order")


def duplicate_slide(prs, source):
    """Add a copy of `source` at the end of the deck, keeping images and notes.

    The subtlety: deep-copied shapes still carry the SOURCE part's r:embed ids. So the
    relationships have to be recreated on the new part under the SAME rIds, or every
    picture on the duplicated slide resolves to nothing.
    """
    dest = prs.slides.add_slide(source.slide_layout)

    # Mirror the source's relationship table exactly, rIds included - the layout rel too,
    # so clearing what add_slide created cannot orphan the layout. The notesSlide rel is
    # the one to skip: a notes part belongs to a single slide, so the duplicate gets its
    # own, created below.
    rels = dest.part.rels
    rels._rels.clear()
    for rel in source.part.rels.values():
        if rel.reltype.endswith("/notesSlide"):
            continue
        rels._rels[rel.rId] = _Relationship(
            rels._base_uri,
            rel.rId,
            rel.reltype,
            target_mode=RTM.EXTERNAL if rel.is_external else RTM.INTERNAL,
            target=rel.target_ref if rel.is_external else rel.target_part,
        )

    # Replace the whole <p:cSld>, not just the shapes. Copying shapes alone silently drops
    # any slide-level <p:bg> - which cost slide 114 its grey backdrop the first time round,
    # and was invisible until the duplicate was rendered next to the original.
    sld, src_sld = dest._element, source._element
    sld.replace(sld.find(qn("p:cSld")), copy.deepcopy(src_sld.find(qn("p:cSld"))))

    src_ovr = src_sld.find(qn("p:clrMapOvr"))
    dst_ovr = sld.find(qn("p:clrMapOvr"))
    if src_ovr is not None:
        if dst_ovr is not None:
            sld.replace(dst_ovr, copy.deepcopy(src_ovr))
        else:
            sld.append(copy.deepcopy(src_ovr))
    for key, value in src_sld.attrib.items():
        sld.set(key, value)

    if source.has_notes_slide:
        dest.notes_slide.notes_text_frame.text = (
            source.notes_slide.notes_text_frame.text
        )
    return dest


def main() -> None:
    if not SRC.exists():
        sys.exit(f"cannot find {SRC} - run this from the folder that holds the deck")

    path = load_path()
    prs = Presentation(str(SRC))
    original_count = len(prs.slides)

    wanted = [n for _, group in path for n in group]
    first_seen: dict[int, int] = {}
    order: list[int] = []          # 0-based indices into prs.slides, in play order

    # Physically duplicate any slide the day visits more than once.
    dupes = 0
    pairs: list[tuple[int, int, int]] = []      # (source slide no, src idx, dup idx)
    for n in wanted:
        if n not in first_seen:
            first_seen[n] = n - 1
            order.append(n - 1)
        else:
            src_idx = first_seen[n]
            duplicate_slide(prs, prs.slides[src_idx])
            dup_idx = len(prs.slides) - 1
            order.append(dup_idx)
            pairs.append((n, src_idx, dup_idx))
            dupes += 1
            print(f"  duplicated slide {n} -> new part for its second visit")

    # Prove the copy is faithful before anything is written: same shape tree, same
    # background, same notes.
    for n, src_idx, dup_idx in pairs:
        src, dup = prs.slides[src_idx], prs.slides[dup_idx]
        src_xml = copy.deepcopy(src._element.find(qn("p:cSld"))).xml
        dup_xml = copy.deepcopy(dup._element.find(qn("p:cSld"))).xml
        if src_xml != dup_xml:
            sys.exit(f"FAILED: the copy of slide {n} is not identical to the original")
        if src.has_notes_slide and (
            src.notes_slide.notes_text_frame.text
            != dup.notes_slide.notes_text_frame.text
        ):
            sys.exit(f"FAILED: notes did not copy for slide {n}")
    print(f"{dupes} slides duplicated and verified identical, {len(order)} stops total")

    # Rewrite sldIdLst to exactly the play order, and drop the rels of everything else so
    # the unused slide parts (and their images) fall out of the saved package.
    sld_id_lst = prs.slides._sldIdLst
    entries = list(sld_id_lst)
    keep_rids = {entries[i].get(qn("r:id")) for i in order}

    for entry in entries:
        sld_id_lst.remove(entry)
    for i in order:
        sld_id_lst.append(entries[i])

    dropped = 0
    for rId in [r for r in list(prs.part.rels) if r not in keep_rids]:
        rel = prs.part.rels[rId]
        if rel.reltype.endswith("/slide"):
            prs.part.drop_rel(rId)
            dropped += 1
    print(f"dropped {dropped} unused slide parts")

    # The custom show points at slides that no longer exist here; strip it.
    pres = prs.part._element
    for existing in pres.findall(qn("p:custShowLst")):
        pres.remove(existing)
        print("removed the custom show (the order is the deck now)")

    prs.save(str(OUT))

    check = Presentation(str(OUT))
    kept = len(check.slides)
    hidden = sum(1 for s in check.slides if s._element.get("show") == "0")
    print(f"\nwrote {OUT.name}")
    print(f"  {kept} slides (from {original_count}), {hidden} hidden")
    print(f"  {SRC.stat().st_size/1e6:.0f}MB -> {OUT.stat().st_size/1e6:.0f}MB")

    if kept != len(order):
        sys.exit(f"FAILED: expected {len(order)} slides, got {kept}")
    if hidden:
        sys.exit(f"FAILED: {hidden} slides are hidden; this deck must be all-visible")

    # Prove every stop still carries content, and report the new numbering.
    empty = [i + 1 for i, s in enumerate(check.slides) if not len(s.shapes)]
    if empty:
        sys.exit(f"FAILED: slides with no shapes: {empty}")
    print("  every slide has content, none hidden")

    print("\nnew numbering, card by card:")
    pos = 1
    for label, group in path:
        span = f"{pos}" if len(group) == 1 else f"{pos}-{pos + len(group) - 1}"
        print(f"  {span:>7}  {label}   (was {', '.join(map(str, group))})")
        pos += len(group)


if __name__ == "__main__":
    main()
