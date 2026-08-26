#!/usr/bin/env python3
"""
Apply the mechanical Payday edits to "AI in Testing Workshop.pptx".

What this DOES do (all low-risk, format-preserving):
  - text corrections for the stale facts in DECK-REVIEW-PAYDAY.md §1-2
  - appends Payday-specific facilitator notes to key slides
  - unhides slide 132 (E2E test generation from scratch)
  - hides the three blank visible slides (131, 138, 153)

What this does NOT do, on purpose:
  - it does not add new slides. The deck is Google-Slides-authored and visually
    consistent; programmatically inserted slides look off-brand next to it. The five
    new slides are specified in full in DECK-REVIEW-PAYDAY.md §7 — build those by hand.
  - it does not reorder slides. Do that by dragging in Slides/PowerPoint.
  - it does not touch anything on a hidden slide except 132.

Usage:
    pip install python-pptx
    python3 apply-payday-deck-edits.py

Writes a NEW file and never modifies the original:
    "AI in Testing Workshop.pptx"  ->  "AI in Testing Workshop - Payday 26Aug2026.pptx"

Every change is printed. Anything it could not find is printed as MISS — check those,
because a MISS usually means the slide was edited after this script was written.
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    sys.exit("python-pptx is not installed.  pip install python-pptx")

SRC = Path("AI in Testing Workshop.pptx")
DST = Path("AI in Testing Workshop - Payday 26Aug2026.pptx")

# ---------------------------------------------------------------------------
# 1. Text replacements:  slide number -> [(find, replace), ...]
#    Substring match, case-sensitive. Formatting of the first run is kept.
# ---------------------------------------------------------------------------
REPLACEMENTS: dict[int, list[tuple[str, str]]] = {
    # Slack channel is dated 6 May 2026
    2: [("20260506", "20260826")],

    # GPT-4 training-token figure is a museum piece
    30: [
        (
            "OpenAI’s GPT-4 was trained on approximately 13 trillion tokens",
            "Frontier models are trained on trillions of tokens",
        ),
    ],

    # Model/context table — every row stale except Gemini 2.5 Pro
    32: [
        ("GPT-5", "GPT-5.6 Terra"),
        ("256k", "1M"),
        ("Claude 3.7 Sonnet Extended", "Claude Sonnet 5"),
        ("192k", "1M"),
        ("Claude 4.0 Sonnet Extended", "Claude Opus 5"),
        ("Gemini 2.5 Pro", "Gemini 2.5 Pro (newest GA Pro)"),
    ],

    # THE blocker: Payday has Cursor, not Copilot
    70: [
        (
            "Done in Gemini and Copilot (as these tools are known and everyone has access)",
            "Done in Gemini (browser) and Cursor — the tools you have access to",
        ),
    ],

    # They have no Claude
    74: [("Using Claude.ai in this example", "Using Gemini in this example")],

    # Setting now defaults to true — not an "enable this first" step
    83: [
        (
            "Enable/disable at github.copilot.chat.codeGeneration.useInstructionFiles",
            "Auto-detected — no setting needed (useInstructionFiles now defaults to true)",
        ),
    ],

    # chat.promptFiles is no longer a setting
    89: [
        (
            "Enable chat.promptFiles setting",
            "No enablement needed — see chat.promptFilesLocations to add folders",
        ),
    ],

    # Folder and frontmatter name disagree in a teaching example
    92: [(".github/skills/web-testing/SKILL.md", ".github/skills/webapp-testing/SKILL.md")],

    # Terminology is half-migrated: chat modes were renamed custom agents
    94: [("Chat modes are stored in", "Custom agents are stored in")],

    # Gemini CLI no longer serves personal Google accounts (18 Jun 2026)
    128: [
        (
            "Copilot in VSCode (alternatives Cursor, Windsurf)",
            "Cursor · Copilot in VS Code (alternatives: Windsurf)",
        ),
        (
            "Claude Code, GeminiCLI, CodexCLI",
            "CLI agents: Claude Code, Codex CLI, Antigravity CLI "
            "(note: Gemini CLI dropped personal Google accounts 18 Jun 2026 — "
            "needs an API key or Code Assist Standard/Enterprise)",
        ),
    ],
}

# ---------------------------------------------------------------------------
# 2. Facilitator notes to APPEND (never replace) to the notes slide.
# ---------------------------------------------------------------------------
MARK = "\n\n--- PAYDAY 26 AUG 2026 ---\n"

NOTES: dict[int, str] = {
    3: "Replace with NEW-1, the timed agenda (DECK-REVIEW-PAYDAY.md §7). Leave "
       "'one very expensive krona' on it — it seeds the 17:00 reveal all day.",

    4: "Write their answers on the whiteboard and LEAVE THEM UP. You check against them "
       "at 17:45. Expect a version of 'write e2e tests without it taking a week'. Say out "
       "loud that they will have a green suite before the afternoon break.",

    11: "LOAD-BEARING. Checking = confirming what you already believe, mechanisable, AI is "
        "excellent at it. Testing = deciding what is worth believing. The whole afternoon "
        "rests on this line.",

    13: "Payday are at 'assistance' — they prompt in Cursor a few times a month. Today "
        "moves them to 'delegation'. That is a bigger jump than this slide implies. Say so.",

    24: "CUT 24-26 (or keep one). Current models get this riddle right, so the slide "
        "undercuts itself live. If you keep it, reframe honestly: they used to fail this, "
        "and 'plausible is not correct' still stands.",

    28: "Unreadable from the back — ~28 tab-aligned rows. Cut, or show five rows where the "
        "models disagreed.",

    32: "Table updated Aug 2026. Say out loud that any table like this is stale within "
        "weeks — the durable lesson is 'check the vendor docs', not the numbers. Note that "
        "Flash reached GA ahead of Pro in the Gemini line.",

    39: "THE SPINE OF THE DAY. Reframe for payroll: a suite can be 100% green and still pay "
        "everybody the wrong amount. You pay this off at 17:00 with the tax cliff.",

    49: "Visual FoxPro to PostgreSQL is the wrong example for JS/TS developers. Cut and keep "
        "46-48, or re-domain.",

    50: "Your own note says 'no experience with this'. Don't demo it. Keep the vision IDEA "
        "as a Gemini browser exercise — screenshot in, test ideas out, zero setup.",

    51: "CUT. Your own note says 'no experience with this', and voice mode costs time you "
        "need for the zero-to-green block.",

    62: "This lifecycle spine is exactly what Payday asked for — 'each step of the testing "
        "lifecycle and how to apply AI'. Keep it and refer back to it all day. Ask the "
        "question on the slide; it gets them talking early.",

    70: "FIXED: was 'Gemini and Copilot'. Payday has CURSOR, not Copilot. Add the line: "
        "your screen is Copilot in VS Code, their screens are Cursor, and the repo is "
        "configured for both on purpose. See NEW-3.",

    73: "KEEP AND PROMOTE — this slide already says Cursor Plan Mode. It is one of the few "
        "slides that is already correct for this audience.",

    74: "FIXED tool reference. Still re-domain the example live: not a school grading "
        "engine but a salary-run change — 400 companies, bank + RSK integrations, "
        "historical payslips, banded tax formulas.",

    76: "Promote this. fixtures/icelandic-test-data.json in the workshop repo is the "
        "concrete artifact - verified kennitalas, ISK boundary values, is-IS number and "
        "date formats. "
        "Five minutes to set up, permanent payoff. Exactly their profile.",

    80: "START OF THE BLOCK THEY SIGNED UP FOR. Protect its time. Next: unhidden slide 132.",

    92: "FIXED folder/name mismatch. Show .agents/skills/ too — both Cursor and Copilot "
        "read that path, so shared knowledge is maintained once.",

    94: "FIXED: 'chat modes' -> 'custom agents'. Cursor equivalent is .cursor/agents/; both "
        "tools also read .claude/agents/ as a COMPAT path, but native paths are what this "
        "repo ships — and the keys differ: Cursor uses readonly:, Copilot uses tools:. "
        "One file cannot enforce both boundaries.",

    97: "Correct as written — subagents do run in parallel now. If you unhide slide 213, "
        "fix it first: 'always the vanilla agent' is no longer true and it contradicts 214.",

    101: "HANDOUT, not a slide. Excellent content, impossible to read on a projector. Print "
         "it; show a five-row version.",

    104: "Good slide. The 'no craftsman carries all the tools' line lands. Note this content "
         "repeats on hidden 221/227/245 — only this one is visible, so no action.",

    106: "Exercise 3 assumes GitHub coding agent access. Payday has none — cut it or rewrite "
         "for Cursor cloud agents. Exercise 2 (Playwright MCP traversal, test CASES first, "
         "review, THEN generate) is the keeper: that review step is the habit to build.",

    109: "Their headline ask was UI/e2e automation. This is the answer. npx playwright "
         "run-test-mcp-server for the test-runner MCP; npx @playwright/mcp@latest for pure "
         "browser automation. Both verified on 1.62.1, neither documented on playwright.dev "
         "— run them once yourself first.",

    111: "DEPRECATED FORMAT on this slide and 112/113: .chatmode.md was renamed .agent.md. "
         "Use npx playwright init-agents --loop=copilot (also: claude, codex, opencode, "
         "vscode, vscode-legacy) to generate the official definitions. Verified on 1.62.1.",

    112: "See notes on 111 — same .chatmode.md correction.",

    113: "See notes on 111 — same .chatmode.md correction. The healer is where the moral "
         "hazard talk goes: a failing test means the test is stale OR the app is broken, "
         "and healing the wrong one leaves you green and lying.",

    118: "EMPTY BODY — everything is in these notes. Either write the body or fold this into "
         "120. Do not present an empty frame.",

    126: "THE GATES. Do not cut this. Errors compound: a vague requirement becomes a wrong "
         "plan becomes fifty wrong tests. Ask the room which gate they would skip, and "
         "remember the answers — you come back to it after the reveal.",

    130: "Feedback form — two minutes now, while the pain is fresh.",

    131: "HIDDEN by this script: the slide was completely blank (zero shapes).",

    132: "UNHIDDEN by this script. THE most relevant slide in the deck for Payday — their "
         "UI/e2e automation is at zero. The prompt in these notes walks an agent through "
         "initialising Playwright, a verification spec, .gitignore and an instructions "
         "file. Build this out into two or three visible slides and put it in the "
         "'zero to a green suite' block after slide 80.",

    136: "Title only, no content. Ask Payday on Tuesday whether they have any Selenium. If "
         "not, cut. If yes, it is an excellent agentic task — mechanical, verifiable, "
         "high-volume, with tests as the referee.",

    137: "PROMOTE. Payday reported NO API automation and they integrate with Icelandic banks "
         "and RSK. Even if you only mention it, this is the obvious second workshop.",

    138: "HIDDEN by this script: the slide was completely blank (zero shapes).",

    139: "CUT 139-143 (the RPI / intentional-compaction section). Genuinely valuable, "
         "genuinely too advanced for a team writing their first e2e test, and it eats 20 "
         "minutes the zero-to-green block needs. Offer it as a follow-up session — it is a "
         "good hook for a second workshop.",

    144: "KEEP, moved to the wrap-up. 'Bad plans come from incomplete or wrong "
         "documentation' is the context-rot lesson, and it is the one habit that pays "
         "forward for the whole team.",

    146: "Timeline stops around Aug 2025 — a year short. Extend it or cut it; a timeline "
         "that ends before the present makes the whole deck feel dated.",

    153: "HIDDEN by this script: the slide was completely blank (zero shapes).",
}

# ---------------------------------------------------------------------------
# 3. Visibility changes
# ---------------------------------------------------------------------------
UNHIDE = [132]           # E2E test generation from scratch — Payday's headline need
HIDE = [131, 138, 153]   # verified blank: zero shapes, no text, no image


# ---------------------------------------------------------------------------
def replace_in_shape(shape, find: str, repl: str) -> int:
    """Replace `find` with `repl`, keeping the first run's formatting."""
    hits = 0
    frames = []
    if shape.has_text_frame:
        frames.append(shape.text_frame)
    if shape.has_table:
        frames.extend(c.text_frame for r in shape.table.rows for c in r.cells)

    for tf in frames:
        for para in tf.paragraphs:
            runs = para.runs
            if not runs:
                continue
            # Fast path: contained in a single run.
            single = False
            for run in runs:
                if find in run.text:
                    run.text = run.text.replace(find, repl)
                    hits += 1
                    single = True
            if single:
                continue
            # Slow path: spans runs. Collapse into the first run.
            joined = "".join(r.text for r in runs)
            if find in joined:
                runs[0].text = joined.replace(find, repl)
                for run in runs[1:]:
                    run.text = ""
                hits += 1
    return hits


def append_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    existing = tf.text or ""
    if MARK.strip() in existing:
        # Already stamped — replace everything after the marker.
        existing = existing.split(MARK.strip())[0].rstrip()
    tf.text = f"{existing}{MARK}{text}" if existing else f"{MARK.strip()}\n{text}"


def set_hidden(slide, hidden: bool) -> None:
    el = slide._element
    if hidden:
        el.set("show", "0")
    elif el.get("show") is not None:
        del el.attrib["show"]


def main() -> None:
    if not SRC.exists():
        sys.exit(f"Cannot find {SRC!s} — run this from the folder that contains it.")

    shutil.copy2(SRC, DST)
    prs = Presentation(str(DST))
    slides = list(prs.slides)
    total = len(slides)
    print(f"Opened copy: {DST}  ({total} slides)\n")

    print("== TEXT REPLACEMENTS ==")
    misses = 0
    for num, pairs in sorted(REPLACEMENTS.items()):
        if num > total:
            print(f"  MISS slide {num}: out of range")
            misses += 1
            continue
        slide = slides[num - 1]
        for find, repl in pairs:
            hits = sum(replace_in_shape(sh, find, repl) for sh in slide.shapes)
            if hits:
                print(f"  slide {num:3}  x{hits}  {find[:52]!r} -> {repl[:52]!r}")
            else:
                print(f"  MISS slide {num:3}  {find[:70]!r}")
                misses += 1

    print("\n== VISIBILITY ==")
    for num in UNHIDE:
        set_hidden(slides[num - 1], False)
        print(f"  slide {num:3}  UNHIDDEN")
    for num in HIDE:
        shape_count = len(slides[num - 1].shapes)
        set_hidden(slides[num - 1], True)
        flag = "" if shape_count == 0 else f"  !! WARNING: {shape_count} shapes, not blank — check this"
        print(f"  slide {num:3}  HIDDEN (shapes={shape_count}){flag}")

    print("\n== FACILITATOR NOTES ==")
    for num, text in sorted(NOTES.items()):
        if num > total:
            print(f"  MISS slide {num}: out of range")
            misses += 1
            continue
        append_notes(slides[num - 1], text)
        print(f"  slide {num:3}  notes appended ({len(text)} chars)")

    prs.save(str(DST))
    print(f"\nSaved: {DST}")
    print(f"Original untouched: {SRC}")
    if misses:
        print(f"\n{misses} MISS(es) — those strings were not found. The slide text was "
              f"probably edited after this script was written; fix by hand.")
    print(
        "\nStill to do BY HAND (see DECK-REVIEW-PAYDAY.md):\n"
        "  - build the five NEW slides in §7 (agenda, the app, two-dialect table,\n"
        "    the tax cliff, commitments)\n"
        "  - turn slide 101 into a printed handout\n"
        "  - reorder: put slide 132 into the block after slide 80\n"
        "  - decide on 24-26 (age riddle), 49 (FoxPro), 51 (voice mode), 136 (Selenium)\n"
        "  - keep 151, 152, 160, 167, 269, 270 HIDDEN — internal or vendor content\n"
    )


if __name__ == "__main__":
    main()
