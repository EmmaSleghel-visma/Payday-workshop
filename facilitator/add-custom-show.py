#!/usr/bin/env python3
"""
Add a PowerPoint Custom Show to the FINAL deck, so the day is one forward path.

The running order is NOT the deck order. Reconstructed from the cards in
run-of-show-extended.html, the afternoon jumps backwards five times:

    block 5   81 -> 83 -> 84 -> 82          (back to 82, the moved slide)
    block 6   86 -> 84                       (back to the two-editor matrix)
    block 7   113-117 -> 128-130 -> 109-111  (back twice)
              -> 113, 114                    (back again, for the exploratory demo)
    block 8   117                            (back once more, for sabotage-and-heal)
    block 9   125-127 -> 149, 150 -> 151 -> 136

Hunting for a slide number in front of a room is the single most avoidable way to lose
the thread, so this writes a Custom Show: a named, ordered subset of the deck. In
PowerPoint it lives under **Slide Show > Custom Slide Show**, and playing it means one
key, forward, all day. Repeats are allowed and intentional - 113/114 and 117 each appear
twice because two different demos use them.

    python3 add-custom-show.py

Reads   : AI in Testing Workshop - Payday 26Aug2026 - FINAL.pptx
Writes  : the same file, in place, adding only <p:custShowLst>. No slide is moved,
          renumbered, hidden or unhidden, so every number in the playbook still holds.
          Re-running replaces the show rather than adding a second one.
"""

from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

DECK = Path("AI in Testing Workshop - Payday 26Aug2026 - FINAL.pptx")
SHOW_NAME = "Payday 26 Aug - run of show"

# One entry per card in the playbook, in the order the cards appear.
#
# The playbook's slide chips name the *anchor* of each card, not everything shown during
# it - so a path built from the chips alone dropped 23 slides the ledger marks USE, most
# of block 6. The COVERAGE ASSERTION below is what caught that, and what stops it
# happening again: every visible slide must appear at least once, or the script refuses.
PATH: list[tuple[str, list[int]]] = [
    ("B1 08:30  Round the table",              [1, 2, 4, 5]),
    ("B1 08:40  Why we test",                  [6, 7, 8, 9, 10, 11, 12, 13, 14]),
    ("B1 08:52  Context is everything",        [15, 16, 18, 19, 20, 21, 22]),
    ("B1 09:00  Critical thinking",            [23, 24, 28]),
    ("B1 09:08  The blindspot story",          [40]),

    ("B2 09:15  Tokens and context",           [30, 31, 32, 33, 34, 35, 36, 37]),
    ("B2 09:25  Managing the window",          [38, 39]),
    ("B2 09:30  Prompting, compressed",        [41, 42, 53, 55, 57, 58, 59]),

    ("B3 09:50  The lifecycle, stage one",     [62, 63, 64]),
    ("B3 09:55  Requirement critique, live",   [65, 66, 67, 68, 69]),
    ("B3 10:05  Their own requirement",        [70, 71, 72, 73]),

    ("B4 10:20  Plans are cheap",              [74, 75]),
    ("B4 10:28  Test data is a decision",      [76, 77]),
    ("B4 10:40  Icelandic edge cases",         [78, 79, 80]),

    ("B5 12:00  Tour the target",              [81, 83]),
    ("B5 12:07  The context experiment",       [84]),
    ("B5 12:17  Playwright from nothing",      [82]),
    ("B5 12:25  Hand over to them",            [85]),

    ("B6 13:00  The map, two dialects",        [86, 84]),
    ("B6 13:03  Where instructions live",      [87, 88, 89, 90, 91, 92]),
    ("B6 13:08  Stored prompts",               [93]),
    ("B6 13:15  Skills",                       [94, 95, 96, 97]),
    ("B6 13:22  Subagents and boundaries",     [98, 99, 100, 101, 102]),
    ("B6 13:28  Scope, and examples over rules", [103, 104]),
    ("B6 13:32  MCP, and what it costs",       [106, 107, 108]),

    ("B7 13:50  The pipeline and the gates",   [113, 114, 115, 116, 117, 118,
                                                128, 129, 130, 133]),
    ("B7 13:55  Full pipeline on Clear paid",  [109, 110, 111, 112]),
    ("B7 14:15  Exploratory with a browser",   [113, 114]),
    ("B7 14:20  THE REVEAL",                   [131, 132]),

    ("B8 14:30  Sabotage and heal",            [117]),
    ("B8 14:36  Stages IV and V",              [119, 120, 121, 122, 123, 124]),

    ("B9 14:45  What we saw today",            [125, 126, 127]),
    ("B9 14:49  Tools, and where to read more", [134, 135, 140]),
    ("B9 14:52  Why it fails, and what next",  [149, 150]),
    ("B9 14:55  Commitments, then feedback",   [151, 136]),
]

# Deliberately NOT in the path. Anything else missing is an error, not a choice.
#
# These 23 are the CUT column of the ledger: still in the file, never shown. Because the
# custom show simply doesn't list them, you cannot land on one by accident - which is the
# quiet benefit of running the show rather than the deck.
EXCLUDED_ON_PURPOSE = {
    105: "printed handout, not a slide",

    25: "age riddle, 1 of 3 - current models get it right, so it undercuts itself",
    26: "age riddle, 2 of 3",
    27: "age riddle, 3 of 3",
    29: "28-row tab-aligned matrix, unreadable from the back; tiers are dated",

    43: "Chain of Thought - cut for time, not for being wrong",
    44: "Chain of Verification - same",
    46: "role playing",
    47: "zero-shot prompting",
    48: "one-shot prompting",
    49: "few-shot prompting",
    50: "X-shot prompting - the Visual FoxPro example is wrong for a JS/TS team",
    51: "Copilot vision - your own notes say no experience with this",
    52: "Copilot Voice Mode - same, and they have no Copilot",
    54: "duplicate of 'Example on outlining a prompt'",
    56: "duplicate of 'Example on outlining a prompt'",

    141: "Selenium to Playwright conversion - not their problem, they have no suite yet",
    142: "API tests from specs - they have no API tests; out of scope for today",
    144: "RPI section 1 of 5 - good material, too advanced for a first e2e suite",
    145: "RPI section 2 of 5 - offer 144-148 as a follow-up session",
    146: "RPI section 3 of 5",
    147: "RPI section 4 of 5",
    148: "RPI section 5 of 5",
    152: "timeline stops a year short of an Aug 2026 workshop",
}


def main() -> None:
    if not DECK.exists():
        sys.exit(f"cannot find {DECK} - run this from the folder that holds the deck")

    backup = DECK.with_suffix(".pptx.before-custom-show")
    if not backup.exists():
        shutil.copy2(DECK, backup)
        print(f"backup written: {backup.name}")

    prs = Presentation(str(DECK))
    slides = list(prs.slides)
    total = len(slides)

    wanted = [n for _, group in PATH for n in group]

    # Every slide on the path must exist and be visible: a hidden slide is skipped even
    # inside a custom show, which would silently drop it from the day.
    problems = []
    for n in wanted:
        if not 1 <= n <= total:
            problems.append(f"slide {n} is outside 1-{total}")
        elif slides[n - 1]._element.get("show") == "0":
            problems.append(f"slide {n} is HIDDEN and would be skipped")
    # COVERAGE ASSERTION. Every visible slide must be on the path, or the custom show
    # would silently drop a slide someone decided was worth keeping. The only permitted
    # absences are the ones named above.
    visible = {i + 1 for i, s in enumerate(slides) if s._element.get("show") != "0"}
    dropped = sorted(visible - set(wanted) - set(EXCLUDED_ON_PURPOSE))
    if dropped:
        problems.append(
            f"{len(dropped)} visible slides are not on the path: "
            + ", ".join(map(str, dropped))
            + "\n     add them to PATH, hide them, or list them in EXCLUDED_ON_PURPOSE"
        )

    if problems:
        sys.exit("refusing to write the show:\n  " + "\n  ".join(problems))

    print(f"coverage: all {len(visible)} visible slides accounted for "
          f"({len(visible) - len(EXCLUDED_ON_PURPOSE)} on the path, "
          f"{len(EXCLUDED_ON_PURPOSE)} excluded on purpose)")

    # Map slide position -> the relationship id presentation.xml uses for it.
    sld_id_lst = prs.slides._sldIdLst
    rids = [entry.get(qn("r:id")) for entry in sld_id_lst]

    pres = prs.part._element
    for existing in pres.findall(qn("p:custShowLst")):
        pres.remove(existing)
        print("replaced the existing custom show")

    cust_show_lst = pres.makeelement(qn("p:custShowLst"), {})
    cust_show = pres.makeelement(qn("p:custShow"), {"name": SHOW_NAME, "id": "0"})
    sld_lst = pres.makeelement(qn("p:sldLst"), {})
    for n in wanted:
        sld_lst.append(pres.makeelement(qn("p:sld"), {qn("r:id"): rids[n - 1]}))
    cust_show.append(sld_lst)
    cust_show_lst.append(cust_show)

    # Schema order inside <p:presentation>: custShowLst sits after notesSz and before
    # defaultTextStyle. Anchor on whichever of the later elements is present.
    anchor = None
    for tag in ("p:photoAlbum", "p:custDataLst", "p:kinsoku", "p:defaultTextStyle",
                "p:modifyVerifier", "p:extLst"):
        found = pres.find(qn(tag))
        if found is not None:
            anchor = found
            break
    if anchor is None:
        pres.append(cust_show_lst)
    else:
        anchor.addprevious(cust_show_lst)

    prs.save(str(DECK))

    check = Presentation(str(DECK))
    shows = check.part._element.findall(qn("p:custShowLst"))
    listed = shows[0].findall(f"{qn('p:custShow')}/{qn('p:sldLst')}/{qn('p:sld')}")
    print(f"\nwrote custom show {SHOW_NAME!r}: {len(listed)} stops")
    print(f"deck unchanged otherwise: {len(check.slides)} slides, "
          f"{sum(1 for s in check.slides if s._element.get('show') != '0')} visible")

    unique = sorted(set(wanted))
    repeats = sorted({n for n in wanted if wanted.count(n) > 1})
    print(f"{len(unique)} distinct slides, {len(wanted) - len(unique)} deliberate repeats "
          f"({', '.join(map(str, repeats))})")

    print("\nthe path, in order:")
    for label, group in PATH:
        print(f"  {label:36} {', '.join(map(str, group))}")

    print(f"\nnot on the path, on purpose ({len(EXCLUDED_ON_PURPOSE)} slides):")
    for n, why in sorted(EXCLUDED_ON_PURPOSE.items()):
        print(f"  {n:4}  {why}")


if __name__ == "__main__":
    main()
